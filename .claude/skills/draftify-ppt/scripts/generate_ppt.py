#!/usr/bin/env python3
"""
Draftify PPT Generator v3.0

Generates planning documents (기획서) in PowerPoint format,
preserving exact template styles including colors, positions, and layouts.

Usage:
    python generate_ppt.py <project_output_dir> [--template <template_path>]
"""

import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    Image = None


# Constants
SCRIPT_DIR = Path(__file__).parent
SKILL_DIR = SCRIPT_DIR.parent
ASSETS_DIR = SKILL_DIR / "assets"
DEFAULT_TEMPLATE = ASSETS_DIR / "ppt_template.pptx"

# Template colors (extracted from template)
COLOR_PRIMARY = RGBColor(0x5E, 0x2B, 0xB8)    # Purple - main accent
COLOR_SECONDARY = RGBColor(0x00, 0x20, 0x60)  # Dark blue - badges
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)

# Template layout names
LAYOUT_COVER = "사용자 지정 레이아웃"
LAYOUT_CONTENT = "4_사용자 지정 레이아웃"
LAYOUT_SECTION = "3_빈화면"
LAYOUT_PROCESS = "6_사용자 지정 레이아웃"
LAYOUT_SCREEN = "7_사용자 지정 레이아웃"
LAYOUT_BLANK = "빈화면"


class MarkdownParser:
    """Parse markdown content into structured data."""

    @staticmethod
    def parse_tables(content: str) -> list[dict]:
        """Extract markdown tables as list of dicts."""
        tables = []
        table_pattern = r'\|(.+)\|\n\|[-:\s|]+\|\n((?:\|.+\|\n?)+)'

        for match in re.finditer(table_pattern, content):
            headers = [h.strip() for h in match.group(1).split('|') if h.strip()]
            rows = []
            for row_line in match.group(2).strip().split('\n'):
                cells = [c.strip() for c in row_line.split('|') if c.strip()]
                if cells:
                    rows.append(cells)
            tables.append({'headers': headers, 'rows': rows})

        return tables

    @staticmethod
    def parse_sections(content: str) -> list[dict]:
        """Parse markdown into sections by headers."""
        sections = []
        current_section = {'level': 0, 'title': '', 'content': ''}

        for line in content.split('\n'):
            if line.startswith('###'):
                if current_section['title']:
                    sections.append(current_section)
                current_section = {'level': 3, 'title': line[3:].strip(), 'content': ''}
            elif line.startswith('##'):
                if current_section['title']:
                    sections.append(current_section)
                current_section = {'level': 2, 'title': line[2:].strip(), 'content': ''}
            elif line.startswith('#'):
                if current_section['title']:
                    sections.append(current_section)
                current_section = {'level': 1, 'title': line[1:].strip(), 'content': ''}
            else:
                current_section['content'] += line + '\n'

        if current_section['title']:
            sections.append(current_section)

        return sections

    @staticmethod
    def strip_markdown(content: str) -> str:
        """Remove markdown formatting for plain text."""
        content = re.sub(r'```[\s\S]*?```', '', content)
        content = re.sub(r'`[^`]+`', '', content)
        content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
        content = re.sub(r'\*([^*]+)\*', r'\1', content)
        content = re.sub(r'^#{1,6}\s+', '', content, flags=re.MULTILINE)
        content = re.sub(r'\|[-:\s|]+\|', '', content)
        return content.strip()


class DraftifyPPTGenerator:
    """Generates planning document PPT preserving template styles."""

    def __init__(self, project_dir: Path, template_path: Path | None = None):
        self.project_dir = Path(project_dir)
        self.template_path = template_path or DEFAULT_TEMPLATE
        self.screenshots_dir = self.project_dir / "screenshots"
        self.analysis_dir = self.project_dir / "analysis"
        self.sections_dir = self.project_dir / "sections"
        # Output filename will be set after loading project name
        self._output_path_base = self.project_dir
        self.output_path = None  # Will be set in load_data()

        self.analyzed_data: dict[str, Any] = {}
        self.sections: dict[str, str] = {}
        self.prs: Presentation | None = None
        self.layouts: dict[str, Any] = {}
        self.parser = MarkdownParser()

    def load_data(self) -> bool:
        """Load analyzed structure and section markdown files."""
        structure_file = self.analysis_dir / "analyzed-structure.json"
        if structure_file.exists():
            with open(structure_file, "r", encoding="utf-8") as f:
                self.analyzed_data = json.load(f)
            print(f"[OK] Loaded: {structure_file.name}")
        else:
            print(f"[WARN] Warning: {structure_file} not found")
            self.analyzed_data = {"project": {"name": "Unknown Project"}}

        # Set output filename based on project name: <project-name>-draft-V0.1.pptx
        project_name = self.analyzed_data.get("project", {}).get("name", "")
        if not project_name or project_name == "Unknown Project":
            # Use directory name as fallback
            project_name = self.project_dir.name
        # Clean project name for filename (replace spaces with hyphens)
        clean_name = project_name.replace(" ", "-").lower()
        version = self.analyzed_data.get("project", {}).get("version", "0.1")
        self.output_path = self._output_path_base / f"{clean_name}-draft-V{version}.pptx"

        # Required sections (warning if missing)
        required_files = {
            "glossary": "05-glossary.md",
            "policy": "06-policy-definition.md",
            "process": "07-process-flow.md",
            "screen": "08-screen-definition.md",
        }

        # Optional sections (use fallback if missing)
        optional_files = {
            "cover": "01-cover.md",
            "history": "02-revision-history.md",
            "toc": "03-table-of-contents.md",
            "divider": "04-section-divider.md",
            "references": "09-references.md",
            "eod": "10-eod.md",
        }

        # Load required sections
        for key, filename in required_files.items():
            filepath = self.sections_dir / filename
            if filepath.exists():
                with open(filepath, "r", encoding="utf-8") as f:
                    self.sections[key] = f.read()
                print(f"[OK] Loaded: {filename}")
            else:
                print(f"[WARN] Warning: {filename} not found")
                self.sections[key] = ""

        # Load optional sections (None = use fallback)
        for key, filename in optional_files.items():
            filepath = self.sections_dir / filename
            if filepath.exists():
                with open(filepath, "r", encoding="utf-8") as f:
                    self.sections[key] = f.read()
                print(f"[OK] Loaded: {filename}")
            else:
                print(f"[INFO] Optional: {filename} not found (using defaults)")
                self.sections[key] = None

        return True

    def create_presentation(self) -> None:
        """Initialize presentation from template."""
        if self.template_path.exists():
            self.prs = Presentation(str(self.template_path))
            print(f"[OK] Using template: {self.template_path.name}")

            # Cache layouts by name
            for layout in self.prs.slide_layouts:
                self.layouts[layout.name] = layout

            # Keep first 2 slides (cover, history), delete the rest
            while len(self.prs.slides) > 2:
                rId = self.prs.slides._sldIdLst[2].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[2]
            print(f"[OK] Kept cover and history slides from template")

            # Remove all PowerPoint sections (slide grouping)
            self._remove_ppt_sections()
            print(f"[OK] Removed PowerPoint sections")
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            print("[WARN] Template not found, using blank presentation")

    def _remove_ppt_sections(self) -> None:
        """Remove all PowerPoint sections (slide grouping) from the presentation."""
        pres_elm = self.prs._element
        # Find and remove sectionLst element (PowerPoint 2010+ sections)
        ns_p14 = '{http://schemas.microsoft.com/office/powerpoint/2010/main}'
        for section_lst in pres_elm.findall(f'.//{ns_p14}sectionLst'):
            section_lst.getparent().remove(section_lst)

    def _get_layout(self, name: str):
        """Get layout by name with fallback."""
        if name in self.layouts:
            return self.layouts[name]
        # Fallback to first available
        return list(self.layouts.values())[0] if self.layouts else self.prs.slide_layouts[0]

    def _add_text_box(self, slide, left: float, top: float, width: float, height: float,
                      text: str, font_size: int = 12, bold: bool = False,
                      font_color: RGBColor = None, fill_color: RGBColor = None,
                      align: PP_ALIGN = PP_ALIGN.LEFT):
        """Add a styled text box."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Remove auto-fit to allow proper sizing
        tf.auto_size = None

        p = tf.paragraphs[0]
        p.alignment = align

        # Create run explicitly to ensure font settings apply
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold

        if font_color:
            run.font.color.rgb = font_color

        if fill_color:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = fill_color

        return txBox

    def _add_rounded_rect(self, slide, left: float, top: float, width: float, height: float,
                          text: str = "", fill_color: RGBColor = COLOR_PRIMARY,
                          font_size: int = 18, font_color: RGBColor = COLOR_WHITE):
        """Add a rounded rectangle with text (template style section bar)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()  # No border

        if text and shape.has_text_frame:
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Create run explicitly for proper font settings
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            run.font.bold = True

        return shape

    def _add_number_badge(self, slide, left: float, top: float, number: str,
                          size: float = 0.39, font_size: int = 14,
                          fill_color: RGBColor = None, font_color: RGBColor = None):
        """Add a small numbered badge (template style - rectangle with gray fill, no border)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(size), Inches(size)
        )
        # Fill color (default: light gray like template)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color if fill_color else RGBColor(0xF2, 0xF2, 0xF2)
        # No border
        shape.line.fill.background()

        if shape.has_text_frame:
            tf = shape.text_frame
            # Match template text frame settings
            tf.word_wrap = False  # Disable text wrapping (template setting)

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Create run - NOT bold for default style (template uses regular weight)
            run = p.add_run()
            run.text = number
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color if font_color else COLOR_BLACK
            # Only bold if custom fill_color is provided (for section badges)
            run.font.bold = True if fill_color else False

        return shape

    def _add_line(self, slide, start_x: float, start_y: float, end_x: float, end_y: float,
                  color: RGBColor = COLOR_GRAY, width: float = 1.0):
        """Add a line shape."""
        connector = slide.shapes.add_connector(
            1,  # Straight connector
            Inches(start_x), Inches(start_y),
            Inches(end_x), Inches(end_y)
        )
        connector.line.color.rgb = color
        connector.line.width = Pt(width)
        return connector

    def _create_table(self, slide, left: float, top: float, width: float, height: float,
                      headers: list, rows: list):
        """Create a styled table matching template."""
        cols = len(headers)
        row_count = len(rows) + 1

        table_shape = slide.shapes.add_table(
            row_count, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        table = table_shape.table

        col_width = Inches(width / cols)
        for i in range(cols):
            table.columns[i].width = col_width

        # Header row styling
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            # Header background (purple)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_PRIMARY
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(10)
                p.font.color.rgb = COLOR_WHITE

        # Data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < cols:
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_text)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(9)

        return table

    # === Slide Generation Methods ===

    def _parse_cover_markdown(self, content: str) -> dict:
        """Parse cover metadata from markdown content."""
        result = {"name": None, "version": None, "date": None, "author": None}
        for line in content.split('\n'):
            line = line.strip()
            # Match: - **Key**: Value or **Key**: Value
            match = re.match(r'^-?\s*\*\*([^*]+)\*\*:\s*(.+)', line)
            if match:
                key = match.group(1).strip().lower()
                value = match.group(2).strip()
                if '프로젝트' in key or 'project' in key or '이름' in key or 'name' in key:
                    result["name"] = value
                elif '버전' in key or 'version' in key:
                    result["version"] = value
                elif '작성일' in key or 'date' in key or '날짜' in key:
                    result["date"] = value
                elif '작성자' in key or 'author' in key:
                    result["author"] = value
            # Match H1 title: # Project Name 기획서
            elif line.startswith('# '):
                title = line[2:].strip()
                # Remove "기획서" suffix if present
                if title.endswith('기획서'):
                    title = title[:-3].strip()
                if not result["name"]:
                    result["name"] = title
        return result

    def add_cover_slide(self) -> None:
        """Update existing cover slide from template - only change title and date."""
        # Use the first slide from template (already preserved in create_presentation)
        slide = self.prs.slides[0]

        # Try to use 01-cover.md first, fallback to analyzed_data
        cover_data = {}
        if self.sections.get("cover"):
            cover_data = self._parse_cover_markdown(self.sections["cover"])

        project_name = cover_data.get("name") or self.analyzed_data.get("project", {}).get("name", "기획서")
        version = cover_data.get("version") or self.analyzed_data.get("project", {}).get("version", "1.0")
        date_str = cover_data.get("date") or datetime.now().strftime("%Y-%m-%d")

        # Find and update existing shapes by name
        for shape in slide.shapes:
            if shape.name == "제목 3" and shape.has_text_frame:
                # Update main title - clear all runs, set text in first run only
                tf = shape.text_frame
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    # Template has multiple runs: "Draftify" + " " + "기획서"
                    # Clear all runs first, then set first run
                    for i, run in enumerate(p.runs):
                        if i == 0:
                            run.text = f"{project_name} 기획서"
                        else:
                            run.text = ""

            elif shape.name == "부제목 4" and shape.has_text_frame:
                # Update date in the subtitle text box (3 lines format)
                tf = shape.text_frame
                for p in tf.paragraphs:
                    if "Last update" in p.text:
                        # Template: "Last update " + ": YYYY-MM-DD"
                        for i, run in enumerate(p.runs):
                            if i == 0:
                                run.text = f"Last update : {date_str}"
                            else:
                                run.text = ""
                    elif "Version" in p.text:
                        if p.runs:
                            p.runs[0].text = f"Version: {version}"

        print("[OK] Updated: Cover slide (template preserved)")

    def _parse_history_markdown(self, content: str) -> list[list[str]]:
        """Parse revision history table from markdown content.
        Returns list of rows, each row is [version, date, changes, author, note]."""
        tables = self.parser.parse_tables(content)
        if tables and tables[0].get('rows'):
            return tables[0]['rows']
        return []

    def add_history_slide(self) -> None:
        """Update existing history slide from template - only update table rows."""
        # Use the second slide from template (already preserved in create_presentation)
        slide = self.prs.slides[1]

        # Try to use 02-revision-history.md first
        history_rows = []
        if self.sections.get("history"):
            history_rows = self._parse_history_markdown(self.sections["history"])

        # Fallback to default if no history from markdown
        if not history_rows:
            version = self.analyzed_data.get("project", {}).get("version", "1.0")
            author = self.analyzed_data.get("project", {}).get("author", "Draftify")
            date_str = datetime.now().strftime("%Y.%m.%d")
            history_rows = [[version, date_str, "기획서 자동 생성", author, ""]]

        # Find the table and update rows
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Row 0 is header, update data rows starting from Row 1
                for ri, row_data in enumerate(history_rows):
                    table_row_idx = ri + 1  # Skip header row
                    if table_row_idx < len(table.rows):
                        row = table.rows[table_row_idx]
                        # Columns: ver, 변경일, 변경내용, 작성자, 비고
                        for ci, cell in enumerate(row.cells):
                            if ci < len(row_data):
                                self._update_cell_text(cell, row_data[ci])
                            else:
                                self._update_cell_text(cell, "")
                break

        print("[OK] Updated: History slide (template preserved)")

    def _update_cell_text(self, cell, new_text: str) -> None:
        """Update cell text while preserving formatting."""
        tf = cell.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            if p.runs:
                # Clear all runs except first, set new text in first run
                for i, run in enumerate(p.runs):
                    if i == 0:
                        run.text = new_text
                    else:
                        run.text = ""
            else:
                # No runs exist, set paragraph text (will lose formatting)
                p.text = new_text
        else:
            cell.text = new_text

    def _extract_section_headers(self, section_key: str) -> list[str]:
        """Extract ## level headers from section markdown."""
        content = self.sections.get(section_key, "")
        headers = []
        for line in content.split('\n'):
            # Match ## headers (not ### or deeper)
            if line.startswith('## ') and not line.startswith('### '):
                # Extract just the title part, remove numbering like "1." or "## 1."
                header = line[3:].strip()
                # Remove leading numbers like "1. " or "1) "
                header = re.sub(r'^[\d]+[\.\)]\s*', '', header)
                if header:
                    headers.append(header)
        return headers

    def _parse_toc_markdown(self, content: str) -> list[tuple[str, str, list[str]]]:
        """Parse TOC from markdown content.
        Returns list of (num, title, sub_items)."""
        toc_items = []
        current_item = None
        current_subs = []

        for line in content.split('\n'):
            line_stripped = line.strip()

            # Match numbered items: 1. Title or ## 1. Title
            num_match = re.match(r'^(?:##\s*)?(\d+)\.\s+(.+)', line_stripped)
            if num_match:
                # Save previous item
                if current_item:
                    toc_items.append((current_item[0], current_item[1], current_subs))
                current_item = (f"{int(num_match.group(1)):02d}", num_match.group(2).strip())
                current_subs = []
            # Match sub-items: - item or   - item (indented)
            elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                if current_item:
                    current_subs.append(line_stripped[2:].strip())

        # Save last item
        if current_item:
            toc_items.append((current_item[0], current_item[1], current_subs))

        return toc_items

    def add_contents_slide(self) -> None:
        """Add contents/TOC slide matching template style with sub-items."""
        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_CONTENT))

        project_name = self.analyzed_data.get("project", {}).get("name", "서비스")

        # "CONTENTS" title - 32pt Bold, Purple (from template)
        self._add_text_box(slide, 1.27, 1.24, 2.44, 0.64, "CONTENTS",
                          font_size=32, bold=True, font_color=COLOR_PRIMARY)

        # Project subtitle - 16pt (from template)
        self._add_text_box(slide, 1.33, 1.86, 4.0, 0.37,
                          f"{project_name} 정책 및 프로세스",
                          font_size=16, font_color=COLOR_WHITE,
                          fill_color=COLOR_PRIMARY)

        # Try to use 03-table-of-contents.md first
        toc_items = []
        if self.sections.get("toc"):
            toc_items = self._parse_toc_markdown(self.sections["toc"])

        # Fallback: build TOC from section headers
        if not toc_items:
            policy_subs = self._extract_section_headers("policy")
            process_subs = self._extract_section_headers("process")
            screen_subs = self._extract_section_headers("screen")

            toc_items = [
                ("01", "용어 정의", []),
                ("02", "정책 정의", policy_subs),
                ("03", "프로세스 흐름도", process_subs),
                ("04", "화면상세", screen_subs),
            ]

        # Calculate total height needed
        total_items = len(toc_items)
        for _, _, subs in toc_items:
            total_items += min(len(subs), 6)  # Cap at 6 sub-items per section

        # Layout: 2 columns
        col1_x = 1.5   # Left column x position
        col2_x = 7.0   # Right column x position
        y_start = 2.6
        line_height_main = 0.5
        line_height_sub = 0.32
        badge_size = 0.39

        current_col = 1
        y = y_start

        for num, title, subs in toc_items:
            # Check if we need to switch to column 2
            items_height = line_height_main + (min(len(subs), 6) * line_height_sub)
            if y + items_height > 6.5 and current_col == 1:
                current_col = 2
                y = y_start

            x = col1_x if current_col == 1 else col2_x

            # Number badge - 14pt
            self._add_number_badge(slide, x, y, num, size=badge_size)
            # Main title - 14pt
            self._add_text_box(slide, x + 0.5, y, 3.0, 0.4, title, font_size=14)
            y += line_height_main

            # Sub-items - 12pt, indented
            for i, sub in enumerate(subs[:6]):  # Max 6 sub-items
                # Truncate long names
                display_sub = sub[:25] + "..." if len(sub) > 25 else sub
                self._add_text_box(slide, x + 0.6, y, 3.5, 0.3,
                                  f"· {display_sub}", font_size=12, font_color=COLOR_GRAY)
                y += line_height_sub

            # Add spacing between sections
            y += 0.15

        print("[OK] Added: Contents slide")

    def _parse_divider_markdown(self, content: str) -> dict[str, str]:
        """Parse section divider titles from markdown content.
        Returns dict mapping section number to title, e.g., {"01": "용어 정의"}."""
        titles = {}
        for line in content.split('\n'):
            line = line.strip()
            # Match: ## 01 Title or ## 01. Title
            match = re.match(r'^##\s*(\d+)\.?\s+(.+)', line)
            if match:
                num = f"{int(match.group(1)):02d}"
                titles[num] = match.group(2).strip()
        return titles

    def _get_divider_title(self, section_num: str, default_title: str) -> str:
        """Get section divider title from 04-section-divider.md or use default."""
        if not hasattr(self, '_divider_titles'):
            self._divider_titles = {}
            if self.sections.get("divider"):
                self._divider_titles = self._parse_divider_markdown(self.sections["divider"])
        return self._divider_titles.get(section_num, default_title)

    def add_section_divider(self, section_num: str, title: str, subtitles: list = None) -> None:
        """Add section divider slide with purple bar (template style)."""
        # Try to get title from 04-section-divider.md
        title = self._get_divider_title(section_num, title)

        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_SECTION))

        # Main section bar (purple rounded rectangle)
        # Position from template: 2.74, 2.46, size: 7.86x1.07
        # Template: 32pt, number in cyan (#08D2D9), title in white
        self._add_section_title_bar(
            slide, 2.74, 2.46, 7.86, 1.07,
            section_num, title
        )

        # Subtitles with number badges (template: 16pt)
        if subtitles:
            y_start = 4.16  # From template position
            for i, subtitle in enumerate(subtitles[:6], 1):
                y = y_start + (i - 1) * 0.52
                self._add_number_badge(slide, 5.93, y, f"{i:02d}", size=0.31)
                self._add_text_box(slide, 6.30, y - 0.03, 2.95, 0.37,
                                  subtitle, font_size=16)

        print(f"[OK] Added: Section divider - {title}")

    def _add_section_title_bar(self, slide, left: float, top: float,
                                width: float, height: float,
                                section_num: str, title: str) -> None:
        """Add section title bar with colored number (template style)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # Template uses rectangle, not rounded
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PRIMARY
        shape.line.fill.background()

        if shape.has_text_frame:
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Number in cyan (#08D2D9)
            run_num = p.add_run()
            run_num.text = f"{section_num} "
            run_num.font.size = Pt(32)
            run_num.font.color.rgb = RGBColor(0x08, 0xD2, 0xD9)  # Cyan
            run_num.font.bold = True

            # Title in white
            run_title = p.add_run()
            run_title.text = title
            run_title.font.size = Pt(32)
            run_title.font.color.rgb = COLOR_WHITE
            run_title.font.bold = True

    def add_glossary_slides(self) -> None:
        """Add glossary slides with proper tables - all tables with pagination."""
        # Extract section titles from markdown for divider
        sections = self.parser.parse_sections(self.sections.get("glossary", ""))
        section_titles = [s['title'] for s in sections if s['level'] == 2]
        self.add_section_divider("01", "용어 정의", section_titles[:6])

        content = self.sections.get("glossary", "")
        if not content:
            return

        # Parse all tables with their section context
        glossary_sections = [s for s in sections if s['level'] == 2]

        max_rows_per_page = 10
        slide = None
        y_pos = 7.0  # Force new slide on first iteration
        page_num = 0

        for section in glossary_sections:
            section_title = section['title']
            tables = self.parser.parse_tables(section['content'])

            for table_data in tables:
                if not table_data['headers'] or not table_data['rows']:
                    continue

                all_rows = table_data['rows']
                headers = table_data['headers']

                # Split rows into chunks for pagination
                for chunk_start in range(0, len(all_rows), max_rows_per_page):
                    chunk_rows = all_rows[chunk_start:chunk_start + max_rows_per_page]
                    row_height = 0.32 * (len(chunk_rows) + 1)

                    # Check if we need a new slide
                    if y_pos + row_height > 6.5 or slide is None:
                        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_CONTENT))
                        page_num += 1

                        # Title with section number badge
                        self._add_number_badge(slide, 0.40, 0.56, "1", size=0.28,
                                              fill_color=COLOR_PRIMARY, font_color=COLOR_WHITE)
                        self._add_text_box(slide, 0.75, 0.50, 10.0, 0.40,
                                          f"용어 정의 - {section_title}", font_size=18, bold=True)
                        y_pos = 1.25

                    self._create_table(
                        slide, 0.40, y_pos, 12.80, row_height,
                        headers, chunk_rows
                    )
                    y_pos += row_height + 0.25

        print(f"[OK] Added: Glossary slides ({page_num} pages)")

    def _parse_policy_items(self, content: str) -> list[dict]:
        """Parse individual policy items (### level) from content."""
        policies = []
        current_policy = None

        for line in content.split('\n'):
            if line.startswith('### '):
                if current_policy:
                    policies.append(current_policy)
                # Extract policy ID and title
                title = line[4:].strip()
                policy_id_match = re.match(r'(POL-[A-Z]+-\d+)[:\s]*(.*)', title)
                if policy_id_match:
                    current_policy = {
                        'id': policy_id_match.group(1),
                        'title': policy_id_match.group(2).strip(),
                        'content': ''
                    }
                else:
                    current_policy = {
                        'id': '',
                        'title': title,
                        'content': ''
                    }
            elif current_policy:
                current_policy['content'] += line + '\n'

        if current_policy:
            policies.append(current_policy)

        return policies

    def _parse_policy_content(self, content: str) -> list[tuple[str, str]]:
        """Parse policy content into key-value pairs."""
        items = []
        current_key = None
        current_value = []

        for line in content.split('\n'):
            line = line.strip()
            if line.startswith('- **') and '**:' in line:
                if current_key:
                    items.append((current_key, '\n'.join(current_value).strip()))
                # Extract key and value
                match = re.match(r'-\s*\*\*([^*]+)\*\*:\s*(.*)', line)
                if match:
                    current_key = match.group(1)
                    current_value = [match.group(2)] if match.group(2) else []
            elif line.startswith('  - ') and current_key:
                current_value.append(line[4:])
            elif line and current_key:
                current_value.append(line)

        if current_key:
            items.append((current_key, '\n'.join(current_value).strip()))

        return items

    def _parse_policy_groups(self, content: str) -> list[dict]:
        """Parse policy markdown into groups with their policies."""
        groups = []
        current_group = None
        current_policy = None

        for line in content.split('\n'):
            if line.startswith('## '):
                # Save previous group
                if current_group:
                    if current_policy:
                        current_group['policies'].append(current_policy)
                    groups.append(current_group)
                # New group
                current_group = {
                    'title': line[3:].strip(),
                    'policies': []
                }
                current_policy = None
            elif line.startswith('### ') and current_group:
                # Save previous policy
                if current_policy:
                    current_group['policies'].append(current_policy)
                # New policy
                title = line[4:].strip()
                policy_id_match = re.match(r'(POL-[A-Z]+-\d+)[:\s]*(.*)', title)
                if policy_id_match:
                    current_policy = {
                        'id': policy_id_match.group(1),
                        'title': policy_id_match.group(2).strip(),
                        'content': ''
                    }
                else:
                    current_policy = {
                        'id': '',
                        'title': title,
                        'content': ''
                    }
            elif current_policy:
                current_policy['content'] += line + '\n'

        # Save last items
        if current_group:
            if current_policy:
                current_group['policies'].append(current_policy)
            groups.append(current_group)

        return groups

    def add_policy_slides(self) -> None:
        """Add policy definition slides - all policies with dynamic pagination."""
        content = self.sections.get("policy", "")
        policy_groups = self._parse_policy_groups(content)

        # Extract subtitles for section divider
        subtitles = []
        for pg in policy_groups[:6]:
            title = pg['title'].split('(')[0].strip()
            title = re.sub(r'^[\d]+[\.\)]\s*', '', title)
            subtitles.append(title[:20])
        self.add_section_divider("02", "정책 정의", subtitles)

        page_num = 0
        slide = None
        y_pos = 7.0  # Force new slide

        for group_idx, pg in enumerate(policy_groups, 1):
            group_title = pg['title']
            # Remove leading number
            group_title_clean = re.sub(r'^[\d]+[\.\)]\s*', '', group_title)

            # Get policies from this group
            policies = pg['policies']

            for policy in policies:
                policy_id = policy['id']
                policy_title = policy['title']
                policy_items = self._parse_policy_content(policy['content'])

                # Estimate height needed for this policy
                policy_height = 0.5 + (len(policy_items) * 0.35)

                # Check if we need a new slide
                if y_pos + policy_height > 6.2 or slide is None:
                    slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_CONTENT))
                    page_num += 1

                    # Group title with badge
                    self._add_number_badge(slide, 0.40, 0.56, str(group_idx), size=0.28,
                                          fill_color=COLOR_PRIMARY, font_color=COLOR_WHITE)
                    self._add_text_box(slide, 0.75, 0.50, 10.0, 0.40,
                                      f"{group_idx}. {group_title_clean}", font_size=16, bold=True)
                    y_pos = 1.20

                # Policy ID and title
                policy_header = f"{policy_id}: {policy_title}" if policy_id else policy_title
                self._add_text_box(slide, 0.50, y_pos, 12.0, 0.35,
                                  policy_header, font_size=12, bold=True,
                                  font_color=COLOR_PRIMARY)
                y_pos += 0.38

                # Policy content items
                for key, value in policy_items:
                    # Truncate long values
                    display_value = value[:100] + "..." if len(value) > 100 else value
                    display_value = display_value.replace('\n', ' / ')
                    item_text = f"• {key}: {display_value}"
                    self._add_text_box(slide, 0.60, y_pos, 12.0, 0.30,
                                      item_text, font_size=10)
                    y_pos += 0.30

                y_pos += 0.15  # Spacing between policies

        print(f"[OK] Added: Policy slides ({page_num} pages)")

    def _split_process_content(self, content: str) -> list[dict]:
        """Split process content into subsections (### level) and code blocks."""
        parts = []
        current_part = {'type': 'text', 'title': '', 'content': ''}
        in_code_block = False
        code_block_content = []

        for line in content.split('\n'):
            if line.startswith('```'):
                if in_code_block:
                    # End of code block
                    parts.append({
                        'type': 'diagram',
                        'title': current_part['title'],
                        'content': '\n'.join(code_block_content)
                    })
                    code_block_content = []
                    in_code_block = False
                else:
                    # Start of code block - save current text if any
                    if current_part['content'].strip():
                        parts.append(current_part)
                    in_code_block = True
                    code_block_content = []
            elif in_code_block:
                code_block_content.append(line)
            elif line.startswith('### '):
                if current_part['content'].strip():
                    parts.append(current_part)
                current_part = {
                    'type': 'text',
                    'title': line[4:].strip(),
                    'content': ''
                }
            else:
                current_part['content'] += line + '\n'

        if current_part['content'].strip():
            parts.append(current_part)

        return parts

    def add_process_slides(self) -> None:
        """Add process flow slides - all sections with dynamic pagination."""
        content = self.sections.get("process", "")
        sections = self.parser.parse_sections(content)
        process_sections = [s for s in sections if s['level'] == 2]

        # Extract subtitles
        subtitles = []
        for ps in process_sections[:6]:
            title = ps['title'].split('(')[0].strip()
            title = re.sub(r'^[\d]+[\.\)]\s*', '', title)
            subtitles.append(title[:20])
        self.add_section_divider("03", "프로세스 흐름도", subtitles)

        page_num = 0

        for idx, ps in enumerate(process_sections, 1):
            section_title = ps['title']
            section_title_clean = re.sub(r'^[\d]+[\.\)]\s*', '', section_title)

            # Parse content into parts (subsections, diagrams, tables)
            parts = self._split_process_content(ps['content'])
            tables = self.parser.parse_tables(ps['content'])

            # Create slides for this section
            slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_PROCESS))
            page_num += 1

            # Section title
            self._add_number_badge(slide, 0.40, 0.56, str(idx), size=0.28,
                                  fill_color=COLOR_PRIMARY, font_color=COLOR_WHITE)
            self._add_text_box(slide, 0.75, 0.50, 10.0, 0.40,
                              f"3-{idx}. {section_title_clean[:35]}", font_size=16, bold=True)

            y_pos = 1.20
            max_y = 6.5

            # Add diagrams and text
            for part in parts:
                if part['type'] == 'diagram':
                    diagram_lines = part['content'].split('\n')
                    diagram_height = len(diagram_lines) * 0.18

                    # Check if we need a new slide
                    if y_pos + diagram_height > max_y:
                        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_PROCESS))
                        page_num += 1
                        self._add_text_box(slide, 0.40, 0.50, 10.0, 0.40,
                                          f"3-{idx}. {section_title_clean[:35]} (계속)",
                                          font_size=14, bold=True)
                        y_pos = 1.10

                    # Subsection title if exists
                    if part['title']:
                        self._add_text_box(slide, 0.50, y_pos, 10.0, 0.30,
                                          part['title'], font_size=11, bold=True)
                        y_pos += 0.35

                    # Add diagram (use monospace styling)
                    diagram_text = part['content'][:1500]  # Limit for very long diagrams
                    self._add_text_box(slide, 0.50, y_pos, 12.0, diagram_height,
                                      diagram_text, font_size=8)
                    y_pos += diagram_height + 0.2

                elif part['type'] == 'text' and part['content'].strip():
                    text_content = part['content'].strip()[:500]
                    text_lines = len(text_content.split('\n'))
                    text_height = max(0.3, text_lines * 0.15)

                    if y_pos + text_height > max_y:
                        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_PROCESS))
                        page_num += 1
                        self._add_text_box(slide, 0.40, 0.50, 10.0, 0.40,
                                          f"3-{idx}. {section_title_clean[:35]} (계속)",
                                          font_size=14, bold=True)
                        y_pos = 1.10

                    if part['title']:
                        self._add_text_box(slide, 0.50, y_pos, 10.0, 0.30,
                                          part['title'], font_size=11, bold=True)
                        y_pos += 0.30

                    self._add_text_box(slide, 0.50, y_pos, 12.0, text_height,
                                      text_content, font_size=9)
                    y_pos += text_height + 0.15

            # Add tables if any
            for table_data in tables:
                if table_data['headers'] and table_data['rows']:
                    rows = table_data['rows']
                    row_height = 0.28 * (len(rows) + 1)

                    if y_pos + row_height > max_y:
                        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_PROCESS))
                        page_num += 1
                        self._add_text_box(slide, 0.40, 0.50, 10.0, 0.40,
                                          f"3-{idx}. {section_title_clean[:35]} (계속)",
                                          font_size=14, bold=True)
                        y_pos = 1.10

                    self._create_table(slide, 0.50, y_pos, 12.0, row_height,
                                      table_data['headers'], rows)
                    y_pos += row_height + 0.2

        print(f"[OK] Added: Process flow slides ({page_num} pages)")

    def _parse_screen_from_markdown(self, content: str) -> dict:
        """Parse a single screen section from markdown into structured data."""
        result = {
            'basic_info': [],
            'ui_elements': {'headers': [], 'rows': []},
            'functions': {'headers': [], 'rows': []},
            'related_policies': []
        }

        current_section = None

        for line in content.split('\n'):
            line_stripped = line.strip()

            # Detect section headers
            if line.startswith('### '):
                section_name = line[4:].strip().lower()
                if '기본' in section_name:
                    current_section = 'basic_info'
                elif 'ui' in section_name or '구성' in section_name:
                    current_section = 'ui_elements'
                elif '기능' in section_name:
                    current_section = 'functions'
                elif '정책' in section_name:
                    current_section = 'policies'
                else:
                    current_section = None
                continue

            # Parse basic info (bullet points)
            if current_section == 'basic_info' and line_stripped.startswith('- **'):
                match = re.match(r'-\s*\*\*([^*]+)\*\*:\s*(.*)', line_stripped)
                if match:
                    result['basic_info'].append((match.group(1), match.group(2)))

            # Parse related policies
            if current_section == 'policies' and line_stripped.startswith('- POL-'):
                result['related_policies'].append(line_stripped[2:])

        # Parse tables from content
        tables = self.parser.parse_tables(content)
        for table in tables:
            if not table['headers']:
                continue
            headers_lower = [h.lower() for h in table['headers']]
            # UI elements table
            if any('요소' in h or 'id' in h.lower() for h in table['headers']):
                if not result['ui_elements']['headers']:
                    result['ui_elements'] = table
            # Functions table
            elif any('기능' in h or '액션' in h or '사용자' in h for h in table['headers']):
                if not result['functions']['headers']:
                    result['functions'] = table

        return result

    def _create_screen_description_table(self, slide, left: float, top: float, width: float,
                                          screen_data: dict, screen_id: str, screen_name: str) -> float:
        """Create description table for screen slide (template style).
        All content from markdown included. Font size fixed at 8pt.
        Returns the y position after the table."""
        # Build description rows from all sections
        rows_data = []

        # Section 1: 기본 정보
        rows_data.append(("기본", ""))  # Section header
        for key, value in screen_data['basic_info']:
            # Clean up value (remove backticks, limit length)
            clean_value = value.replace('`', '').strip()
            display_value = clean_value[:50] + "..." if len(clean_value) > 50 else clean_value
            rows_data.append(("", f"{key}: {display_value}"))

        # Section 2: UI 구성 요소
        ui_table = screen_data['ui_elements']
        if ui_table.get('rows'):
            rows_data.append(("UI", ""))  # Section header
            for row in ui_table['rows']:
                if len(row) >= 3:
                    # Format: ID | 유형 | 설명
                    element_text = f"{row[0]} ({row[1]}): {row[2][:30]}"
                    rows_data.append(("", element_text))

        # Section 3: 기능 정의
        func_table = screen_data['functions']
        if func_table.get('rows'):
            rows_data.append(("기능", ""))  # Section header
            for row in func_table['rows']:
                if len(row) >= 2:
                    # Format: 기능명 | 액션
                    func_text = f"{row[0]}: {row[1][:30]}"
                    if len(row) >= 4 and row[3]:
                        func_text += f" -> {row[3]}"
                    rows_data.append(("", func_text))

        # Section 4: 관련 정책
        if screen_data['related_policies']:
            rows_data.append(("정책", ""))  # Section header
            for policy in screen_data['related_policies']:
                rows_data.append(("", policy))

        # Calculate dynamic row height based on content
        # More rows = smaller row height to fit
        total_rows = len(rows_data)
        available_height = 6.8  # Available height in the description area
        row_height_each = min(0.35, available_height / max(total_rows, 1))

        # Create the table
        table_height = row_height_each * total_rows
        table = slide.shapes.add_table(
            total_rows, 2,
            Inches(left), Inches(top), Inches(width), Inches(table_height)
        ).table

        # Set column widths (narrow first col for section labels)
        table.columns[0].width = Inches(0.35)
        table.columns[1].width = Inches(width - 0.35)

        # Fill table data with font size 8pt fixed
        for ri, (label, desc) in enumerate(rows_data):
            # Label cell (section header or empty)
            cell0 = table.cell(ri, 0)
            cell0.text = label
            if label:  # Section header row
                cell0.fill.solid()
                cell0.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
            else:
                cell0.fill.background()
            for para in cell0.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(8)
                    run.font.bold = True if label else False
                    run.font.color.rgb = COLOR_BLACK

            # Description cell
            cell1 = table.cell(ri, 1)
            cell1.text = desc
            if label:  # Section header row
                cell1.fill.solid()
                cell1.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
            else:
                cell1.fill.background()
            for para in cell1.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(8)
                    run.font.bold = True if label else False
                    run.font.color.rgb = COLOR_BLACK

        # Apply dark gray thin borders to all cells
        # Border width: 0.5pt = 6350 EMU, Color: dark gray #666666
        for ri in range(total_rows):
            for ci in range(2):
                cell = table.cell(ri, ci)
                self._set_cell_border_all(cell, 6350, "666666")

        return top + table_height

    def _set_cell_border_all(self, cell, width_emu: int, color_hex: str) -> None:
        """Set all borders for a cell (dark gray thin line)."""
        from lxml import etree

        tc = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))

        for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
            # Remove existing border
            border = tcPr.find(qn(f'a:{border_name}'))
            if border is not None:
                tcPr.remove(border)

            # Create new border
            border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
            border.set('w', str(width_emu))
            border.set('cap', 'flat')
            border.set('cmpd', 'sng')
            border.set('algn', 'ctr')

            # Solid fill with color
            solidFill = etree.SubElement(border, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', color_hex)

            # Solid dash
            prstDash = etree.SubElement(border, qn('a:prstDash'))
            prstDash.set('val', 'solid')

    def _extract_screen_sections_raw(self, content: str) -> list[dict]:
        """Extract screen sections with full content including ### subsections."""
        screens = []
        current_screen = None

        for line in content.split('\n'):
            # Match ## SCR-XXX: Name
            if line.startswith('## ') and 'SCR-' in line:
                # Save previous screen
                if current_screen:
                    screens.append(current_screen)
                # Start new screen
                match = re.search(r'(SCR-\d{3})[:\s]*(.*)', line[3:])
                if match:
                    current_screen = {
                        'id': match.group(1),
                        'name': match.group(2).strip(),
                        'content': ''
                    }
                else:
                    current_screen = None
            elif line.startswith('## ') and current_screen:
                # Another ## section (not SCR-), save current and stop
                screens.append(current_screen)
                current_screen = None
            elif current_screen:
                current_screen['content'] += line + '\n'

        # Save last screen
        if current_screen:
            screens.append(current_screen)

        return screens

    def add_screen_slides(self) -> None:
        """Add screen definition slides with description in right 25% area (template style)."""
        screen_content = self.sections.get("screen", "")

        # Extract screen sections with full content (including ### subsections)
        screen_sections_raw = self._extract_screen_sections_raw(screen_content)

        # Extract screen IDs and names for divider
        subtitles = []
        for s in screen_sections_raw[:6]:
            subtitles.append(f"{s['id']}: {s['name'][:12]}")
        self.add_section_divider("04", "화면상세", subtitles)

        # Screen list summary slide from markdown table (find "화면 목록 요약" section)
        summary_match = re.search(r'##\s*화면\s*목록[^\n]*\n(.*?)(?=\n##\s|$)', screen_content, re.DOTALL)
        if summary_match:
            tables = self.parser.parse_tables(summary_match.group(1))
            if tables and tables[0]['rows']:
                slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_CONTENT))
                self._add_number_badge(slide, 0.40, 0.56, "4", size=0.28,
                                      fill_color=COLOR_PRIMARY, font_color=COLOR_WHITE)
                self._add_text_box(slide, 0.75, 0.50, 4.0, 0.40,
                                  "화면 목록", font_size=18, bold=True)
                table_data = tables[0]
                row_count = min(len(table_data['rows']), 15)
                self._create_table(slide, 0.40, 1.25, 12.80, 0.32 * (row_count + 1),
                                  table_data['headers'], table_data['rows'][:row_count])

        page_num = 0

        # Template layout: right 25% for description
        # Slide width: 13.333 inches, Description area: ~3.3 inches
        desc_left = 10.0   # Start of description area
        desc_width = 3.2   # Width of description area
        screenshot_width = 9.3  # Left area for screenshot (75%)

        # Individual screen slides
        for idx, screen_section in enumerate(screen_sections_raw, 1):
            screen_id = screen_section['id']
            screen_name = screen_section['name']

            # Parse screen content (includes ### subsections)
            screen_data = self._parse_screen_from_markdown(screen_section['content'])

            # === Main Slide: Screenshot + Description ===
            slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_SCREEN))
            page_num += 1

            # Section title at top left
            self._add_text_box(slide, 0.40, 0.10, 9.0, 0.35,
                              "4. 화면상세", font_size=12, bold=True)

            # Screen name with badge
            self._add_number_badge(slide, 0.40, 0.50, str(idx), size=0.28,
                                  fill_color=COLOR_PRIMARY, font_color=COLOR_WHITE)
            self._add_text_box(slide, 0.75, 0.45, 8.0, 0.35,
                              f"{screen_id}: {screen_name}", font_size=13, bold=True)

            # "Description" header bar (template style - light gray)
            desc_header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(desc_left), Inches(0.10), Inches(desc_width), Inches(0.28)
            )
            desc_header.fill.solid()
            desc_header.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            desc_header.line.fill.background()
            if desc_header.has_text_frame:
                tf = desc_header.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = "Description"
                run.font.size = Pt(9)
                run.font.color.rgb = COLOR_BLACK
                run.font.bold = False

            # Screenshot on left side (75% area)
            screenshot_path = self.screenshots_dir / f"{screen_id}.png"
            if screenshot_path.exists():
                try:
                    # Calculate appropriate size within left area
                    slide.shapes.add_picture(
                        str(screenshot_path),
                        Inches(0.4), Inches(1.0),
                        width=Inches(min(screenshot_width - 0.8, 8.5))
                    )
                except Exception as e:
                    print(f"  [WARN] Screenshot error for {screen_id}: {e}")

            # Description table on right side (25% area)
            # All content (basic info, UI elements, functions, policies) included in this table
            self._create_screen_description_table(
                slide, desc_left, 0.42, desc_width,
                screen_data, screen_id, screen_name
            )

        print(f"[OK] Added: Screen definition slides ({page_num} pages)")

    def _set_cell_borders_selective(self, cell, width_emu: int, color_hex: str,
                                      row: int, col: int, total_rows: int, total_cols: int) -> None:
        """Set cell borders selectively (top/bottom/inside only, no outer left/right).

        Matches template XML structure exactly:
        <a:lnT w="6350" cap="flat" cmpd="sng" algn="ctr">
          <a:solidFill><a:srgbClr val="808080"/></a:solidFill>
          <a:prstDash val="solid"/>
          <a:round/>
          <a:headEnd type="none" w="med" len="med"/>
          <a:tailEnd type="none" w="med" len="med"/>
        </a:lnT>
        """
        from pptx.oxml.ns import qn
        from lxml import etree

        tc = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))

        # Determine which borders to show
        borders_to_set = []
        borders_to_set.append('lnT')  # Top border always
        borders_to_set.append('lnB')  # Bottom border always

        # Inside vertical borders (between columns)
        if col < total_cols - 1:
            borders_to_set.append('lnR')  # Right border if not last column
        if col > 0:
            borders_to_set.append('lnL')  # Left border if not first column

        # Set borders with complete XML structure matching template
        for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
            border = tcPr.find(qn(f'a:{border_name}'))

            if border_name in borders_to_set:
                # Remove existing border if present and recreate
                if border is not None:
                    tcPr.remove(border)

                # Create new border element with all attributes
                border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                border.set('w', str(width_emu))
                border.set('cap', 'flat')
                border.set('cmpd', 'sng')
                border.set('algn', 'ctr')

                # solidFill with srgbClr
                solidFill = etree.SubElement(border, qn('a:solidFill'))
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgbClr.set('val', color_hex)

                # prstDash - solid line
                prstDash = etree.SubElement(border, qn('a:prstDash'))
                prstDash.set('val', 'solid')

                # round - rounded line join
                etree.SubElement(border, qn('a:round'))

                # headEnd - line start cap
                headEnd = etree.SubElement(border, qn('a:headEnd'))
                headEnd.set('type', 'none')
                headEnd.set('w', 'med')
                headEnd.set('len', 'med')

                # tailEnd - line end cap
                tailEnd = etree.SubElement(border, qn('a:tailEnd'))
                tailEnd.set('type', 'none')
                tailEnd.set('w', 'med')
                tailEnd.set('len', 'med')
            else:
                # Remove border completely or set to noFill
                if border is not None:
                    tcPr.remove(border)
                # Create empty border with noFill
                border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                border.set('w', '0')
                etree.SubElement(border, qn('a:noFill'))

    def _parse_references_markdown(self, content: str) -> list[str]:
        """Parse references from markdown content."""
        refs = []
        for line in content.split('\n'):
            line = line.strip()
            # Match bullet points: - item or * item
            if line.startswith('- ') or line.startswith('* '):
                refs.append(line[2:].strip())
            # Match numbered items: 1. item
            elif re.match(r'^\d+\.\s+', line):
                refs.append(re.sub(r'^\d+\.\s+', '', line).strip())
        return refs

    def add_reference_slide(self) -> None:
        """Add reference slide."""
        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_CONTENT))

        self._add_text_box(slide, 0.35, 0.55, 3.0, 0.64,
                          "참고 문헌", font_size=24, bold=True)

        # Try to use 09-references.md first, fallback to analyzed_data
        refs = []
        if self.sections.get("references"):
            refs = self._parse_references_markdown(self.sections["references"])

        if not refs:
            refs = self.analyzed_data.get("references", [])

        ref_text = "\n".join([f"• {ref}" for ref in refs]) if refs else "• 자동 생성된 기획서\n• Draftify 크롤링 데이터 기반"

        self._add_text_box(slide, 0.51, 1.5, 12.42, 4.0, ref_text, font_size=12)

        print("[OK] Added: Reference slide")

    def _parse_eod_markdown(self, content: str) -> str:
        """Parse EOD text from markdown content."""
        for line in content.split('\n'):
            line = line.strip()
            # Match H1: # End of Document
            if line.startswith('# '):
                return line[2:].strip()
            # Match plain text (non-empty, non-header)
            if line and not line.startswith('#'):
                return line
        return "End of document"

    def add_eod_slide(self) -> None:
        """Add end of document slide."""
        slide = self.prs.slides.add_slide(self._get_layout(LAYOUT_BLANK))

        # Try to use 10-eod.md first, fallback to default
        eod_text = "End of document"
        if self.sections.get("eod"):
            eod_text = self._parse_eod_markdown(self.sections["eod"])

        # Centered EOD text (from template: 3.34, 3.14)
        self._add_text_box(slide, 3.34, 3.14, 6.82, 0.61,
                          eod_text,
                          font_size=28, bold=True, font_color=COLOR_GRAY,
                          align=PP_ALIGN.CENTER)

        print("[OK] Added: EOD slide")

    def generate(self) -> Path:
        """Generate the complete PPT document."""
        print(f"\n{'='*50}")
        print("Draftify PPT Generator v3.0")
        print(f"{'='*50}\n")

        print("Loading data...")
        self.load_data()

        print("\nCreating presentation...")
        self.create_presentation()

        print("\nGenerating slides...")
        self.add_cover_slide()
        self.add_history_slide()
        self.add_contents_slide()
        self.add_glossary_slides()
        self.add_policy_slides()
        self.add_process_slides()
        self.add_screen_slides()
        self.add_reference_slide()
        self.add_eod_slide()

        print(f"\nSaving to: {self.output_path}")
        self.prs.save(str(self.output_path))

        print(f"\n{'='*50}")
        print(f"[OK] SUCCESS: Generated {self.output_path}")
        print(f"  Total slides: {len(self.prs.slides)}")
        print(f"{'='*50}\n")

        return self.output_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate planning document PPT from analyzed project data"
    )
    parser.add_argument(
        "project_dir",
        help="Path to project output directory (e.g., outputs/my-project/)"
    )
    parser.add_argument(
        "--template",
        help="Path to PPT template file",
        default=None
    )

    args = parser.parse_args()

    project_dir = Path(args.project_dir)
    if not project_dir.exists():
        print(f"Error: Project directory not found: {project_dir}")
        sys.exit(1)

    template_path = Path(args.template) if args.template else None

    generator = DraftifyPPTGenerator(project_dir, template_path)
    generator.generate()


if __name__ == "__main__":
    main()
